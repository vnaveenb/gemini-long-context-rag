"""PPTX text extraction using python-pptx."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from src.logger import get_logger
from src.models.document import (
    DocumentFormat,
    DocumentMetadata,
    ExtractedContent,
    PageContent,
    Section,
)

logger = get_logger(__name__)


def extract_pptx(file_path: str | Path) -> ExtractedContent:
    """Extract text, sections, and metadata from a PPTX file."""
    file_path = Path(file_path)
    logger.info("Extracting PPTX", path=str(file_path))

    errors: list[str] = []

    try:
        prs = Presentation(str(file_path))
    except Exception as exc:
        logger.error("Failed to open PPTX", path=str(file_path), error=str(exc))
        return ExtractedContent(
            filename=file_path.name,
            format=DocumentFormat.PPTX,
            extraction_errors=[f"Failed to open PPTX: {exc}"],
        )

    pages: list[PageContent] = []
    sections: list[Section] = []
    all_text_parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        try:
            slide_texts: list[str] = []
            slide_title = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)

                    # Detect title shape
                    if shape.shape_type and "TITLE" in str(shape.shape_type).upper():
                        slide_title = shape.text_frame.text.strip()

                # Tables in slides
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_texts.append(" | ".join(cells))

            # Speaker notes
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[Speaker Notes] {notes}")

            page_text = "\n".join(slide_texts)
            pages.append(
                PageContent(
                    page_number=slide_num,
                    text=page_text,
                    metadata={"title": slide_title, "has_notes": bool(notes)},
                )
            )
            all_text_parts.append(page_text)

            if slide_title:
                sections.append(
                    Section(
                        title=slide_title,
                        level=1,
                        page_start=slide_num,
                        content=page_text,
                    )
                )
        except Exception as exc:
            errors.append(f"Slide {slide_num}: {exc}")
            logger.warning("Error extracting slide", slide=slide_num, error=str(exc))

    # Fill section page_end
    for i, sec in enumerate(sections):
        sec.page_end = sections[i + 1].page_start if i + 1 < len(sections) else len(pages)

    raw_text = "\n\n".join(all_text_parts)

    metadata = DocumentMetadata(
        page_count=len(pages),
        word_count=len(raw_text.split()),
    )

    content = ExtractedContent(
        filename=file_path.name,
        format=DocumentFormat.PPTX,
        metadata=metadata,
        pages=pages,
        sections=sections,
        raw_text=raw_text,
        extraction_errors=errors,
    )
    logger.info(
        "PPTX extraction complete",
        slides=len(pages),
        sections=len(sections),
        words=metadata.word_count,
    )
    return content
